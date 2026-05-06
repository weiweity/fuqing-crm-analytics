"""
芙清 CRM 客户分析系统 - PPT 导出服务
Week 4 PPT 报告生成
"""

import os
import uuid
from datetime import datetime
from typing import Dict, Any, List
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from backend.config import PROJECT_ROOT, DUCKDB_PATH


# 可用模板
AVAILABLE_TEMPLATES = {
    "default": {
        "name": "默认模板",
        "description": "标准商务风格，适用于常规分析报告",
        "modules": ["cover", "metrics", "segments", "geo", "category", "actions"]
    }
}

# 模块列表
MODULE_LIST = ["cover", "metrics", "segments", "geo", "category", "actions"]


def _ensure_export_dir() -> Path:
    """确保导出目录存在"""
    export_dir = PROJECT_ROOT / "data" / "exports"
    export_dir.mkdir(parents=True, exist_ok=True)
    return export_dir


def get_available_templates() -> Dict[str, Any]:
    """
    获取可用模板列表

    Returns:
        {
            "templates": [{"id": str, "name": str, "description": str}, ...],
            "modules": [str, ...]
        }
    """
    templates = []
    for template_id, template_info in AVAILABLE_TEMPLATES.items():
        templates.append({
            "id": template_id,
            "name": template_info["name"],
            "description": template_info["description"],
            "modules": template_info["modules"]
        })

    return {
        "templates": templates,
        "modules": MODULE_LIST
    }


def _create_cover_slide(prs: 'Presentation', title: str, subtitle: str, date_range: str):
    """创建封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

    # 日期范围
    left = Inches(1)
    top = Inches(5)
    width = Inches(8)
    height = Inches(0.5)
    date_box = slide.shapes.add_textbox(left, top, width, height)
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = date_range
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER


def _create_metrics_slide(prs: 'Presentation', metrics_data: Dict[str, Any]):
    """创建核心指标页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心指标概览"
    p.font.size = Pt(32)
    p.font.bold = True

    # 指标内容
    content = f"""
    GMV: ¥{metrics_data.get('amount', 0):,.2f}
    订单数: {metrics_data.get('order_count', 0):,}
    客单价: ¥{metrics_data.get('avg_order_value', 0):,.2f}

    新客: {metrics_data.get('new_users', 0):,} (¥{metrics_data.get('new_user_amount', 0):,.2f})
    老客: {metrics_data.get('old_users', 0):,} (¥{metrics_data.get('old_user_amount', 0):,.2f})

    会员金额: ¥{metrics_data.get('member_amount', 0):,.2f}
    会员占比: {metrics_data.get('member_ratio', 0):.1f}%

    环比变化: GMV {metrics_data.get('mom_change', {}).get('amount_pct', 0):+.1f}%
    同比变化: GMV {metrics_data.get('yoy_change', {}).get('amount_pct', 0):+.1f}%
    """

    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.strip()
    p.font.size = Pt(18)


def _create_segments_slide(prs: 'Presentation', segments_data: List[Dict[str, Any]]):
    """创建象限分布页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "客户象限分布"
    p.font.size = Pt(32)
    p.font.bold = True

    # 象限内容
    content_lines = []
    for seg in segments_data[:9]:
        content_lines.append(
            f"{seg.get('name', '未知')}: {seg.get('user_count', 0):,} 用户 "
            f"(¥{seg.get('gmv', 0):,.2f})"
        )

    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\n".join(content_lines)
    p.font.size = Pt(18)


def _create_geo_slide(prs: 'Presentation', geo_data: Dict[str, Any]):
    """创建地域分布页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"地域分布 - {geo_data.get('level', '省份')}"
    p.font.size = Pt(32)
    p.font.bold = True

    # 地域内容
    dist = geo_data.get('distribution', [])
    content_lines = [f"总用户: {geo_data.get('total_users', 0):,}", f"总GMV: ¥{geo_data.get('total_gmv', 0):,.2f}", ""]
    for item in dist[:15]:
        content_lines.append(
            f"{item.get('name', '未知')}: {item.get('user_count', 0):,} 用户 "
            f"({item.get('占比', 0):.1f}%)"
        )

    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\n".join(content_lines)
    p.font.size = Pt(16)


def _create_category_slide(prs: 'Presentation', category_data: Dict[str, Any]):
    """创建品类分布页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"品类分布 - {category_data.get('level', 'category')}"
    p.font.size = Pt(32)
    p.font.bold = True

    # 品类内容
    dist = category_data.get('distribution', [])
    content_lines = [f"总用户: {category_data.get('total_users', 0):,}", f"总GMV: ¥{category_data.get('total_gmv', 0):,.2f}", ""]
    for item in dist[:15]:
        content_lines.append(
            f"{item.get('name', '未知')}: {item.get('user_count', 0):,} 用户 "
            f"({item.get('占比', 0):.1f}%)"
        )

    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\n".join(content_lines)
    p.font.size = Pt(16)


def _create_actions_slide(prs: 'Presentation', actions_data: Dict[str, Any]):
    """创建行动建议页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "行动建议"
    p.font.size = Pt(32)
    p.font.bold = True

    # 建议内容
    recommendations = actions_data.get('recommendations', [
        "1. 重点关注高价值用户群体",
        "2. 加强流失用户召回",
        "3. 优化品类结构",
    ])

    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\n".join(recommendations)
    p.font.size = Pt(20)


def generate_ppt_report(
    report_type: str,
    start_date: str,
    end_date: str,
    modules: List[str],
    template: str = "default"
) -> Dict[str, Any]:
    """
    生成 PPT 报告

    Args:
        report_type: 报告类型（如 "weekly", "monthly"）
        start_date: 开始日期 (YYYY-MM-DD)
        end_date: 结束日期 (YYYY-MM-DD)
        modules: 包含的模块列表，如 ["cover", "metrics", "segments", "geo", "category"]
        template: 模板名称，默认 "default"

    Returns:
        {
            "report_id": str,
            "file_name": str,
            "download_url": str
        }
    """
    if not HAS_PPTX:
        return {
            "report_id": None,
            "file_name": None,
            "download_url": None,
            "error": "python-pptx 库未安装，请运行: pip install python-pptx"
        }

    # 验证模板
    if template not in AVAILABLE_TEMPLATES:
        template = "default"

    # 验证模块
    valid_modules = [m for m in modules if m in MODULE_LIST]
    if not valid_modules:
        valid_modules = ["cover", "metrics"]

    # 生成报告 ID
    report_id = str(uuid.uuid4())[:8]
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    file_name = f"CRM_Report_{report_type}_{timestamp}.pptx"

    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 日期范围
    date_range = f"{start_date} 至 {end_date}"

    # 导入数据服务
    from backend.services.metrics_service import get_overview_metrics
    from backend.services.flow_service import get_flow_matrix
    from backend.services.geo_service import get_geo_distribution
    from backend.services.category_service import get_category_distribution

    # 生成各模块幻灯片
    for module in valid_modules:
        if module == "cover":
            _create_cover_slide(prs, "芙清 CRM 客户分析报告", report_type.upper(), date_range)

        elif module == "metrics":
            try:
                metrics_data = get_overview_metrics(start_date, end_date)
                _create_metrics_slide(prs, metrics_data)
            except Exception as e:
                print(f"Warning: Failed to generate metrics slide: {e}")

        elif module == "segments":
            try:
                flow_data = get_flow_matrix(start_date, end_date, lookback_days=90)
                # 提取象限汇总
                segments_summary = []
                for seg in flow_data.get('segments', []):
                    seg_id = seg.get('id')
                    # 获取该象限用户数（简化计算）
                    seg_count = sum(
                        m['count'] for m in flow_data.get('flow_matrix', [])
                        if m.get('from') == seg_id or m.get('to') == seg_id
                    )
                    segments_summary.append({
                        'name': seg.get('name', '未知'),
                        'user_count': seg_count,
                        'gmv': 0
                    })
                _create_segments_slide(prs, segments_summary)
            except Exception as e:
                print(f"Warning: Failed to generate segments slide: {e}")

        elif module == "geo":
            try:
                geo_data = get_geo_distribution(end_date, lookback_days=90, level="省份", top_n=20)
                _create_geo_slide(prs, geo_data)
            except Exception as e:
                print(f"Warning: Failed to generate geo slide: {e}")

        elif module == "category":
            try:
                category_data = get_category_distribution(end_date, lookback_days=90, level="category")
                _create_category_slide(prs, category_data)
            except Exception as e:
                print(f"Warning: Failed to generate category slide: {e}")

        elif module == "actions":
            _create_actions_slide(prs, {"recommendations": [
                "1. 加强高价值客户维护",
                "2. 关注流失风险用户",
                "3. 优化地域投放策略",
                "4. 调整品类结构"
            ]})

    # 保存文件
    export_dir = _ensure_export_dir()
    file_path = export_dir / file_name
    prs.save(str(file_path))

    return {
        "report_id": report_id,
        "file_name": file_name,
        "download_url": f"/data/exports/{file_name}"
    }


if __name__ == "__main__":
    # 测试
    print("=== 可用模板 ===")
    templates = get_available_templates()
    print(f"模板: {[t['id'] for t in templates['templates']]}")
    print(f"模块: {templates['modules']}")

    print("\n=== 生成 PPT 测试 ===")
    result = generate_ppt_report(
        report_type="weekly",
        start_date="2026-03-01",
        end_date="2026-03-19",
        modules=["cover", "metrics", "segments", "geo", "category"],
        template="default"
    )
    print(f"结果: {result}")
